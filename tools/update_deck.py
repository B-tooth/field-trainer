from __future__ import annotations

import argparse
import json
import re
import subprocess
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("Missing dependency: python-pptx. Run tools\\setup-tools.bat once, then try again.")
    raise SystemExit(1)


def slugify(value: str) -> str:
    value = re.sub(r"[^a-zA-Z0-9]+", "_", value.strip().lower())
    return value.strip("_") or "deck"


def find_powerpoint(deck_folder: Path) -> Path:
    files = sorted(deck_folder.glob("*.pptx"))
    if len(files) != 1:
        raise ValueError(f"Expected exactly one .pptx file in {deck_folder}, found {len(files)}.")
    return files[0]


def read_titles(powerpoint: Path) -> list[str]:
    presentation = Presentation(powerpoint)
    titles: list[str] = []
    for number, slide in enumerate(presentation.slides, start=1):
        title = slide.shapes.title
        text = title.text.strip() if title is not None else ""
        if not text:
            raise ValueError(f"Slide {number} has no PowerPoint title placeholder text.")
        titles.append(text)
    return titles


def export_slides(powerpoint: Path, deck_folder: Path, script_folder: Path) -> None:
    export_script = script_folder / "export_slides.ps1"
    command = [
        "powershell.exe", "-NoProfile", "-ExecutionPolicy", "Bypass",
        "-File", str(export_script),
        "-PowerPointPath", str(powerpoint),
        "-OutputFolder", str(deck_folder),
    ]
    completed = subprocess.run(command, text=True, capture_output=True)
    if completed.returncode != 0:
        detail = completed.stderr.strip() or completed.stdout.strip()
        raise RuntimeError(f"PowerPoint slide export failed. {detail}")


def rebuild_index(project_root: Path) -> None:
    decks_folder = project_root / "decks"
    entries = []
    for deck_json in sorted(decks_folder.glob("*/deck.json")):
        data = json.loads(deck_json.read_text(encoding="utf-8"))
        entries.append({"id": data["id"], "name": data["name"], "path": deck_json.relative_to(project_root).as_posix()})
    (decks_folder / "index.json").write_text(json.dumps({"decks": entries}, indent=2) + "\n", encoding="utf-8")


def main() -> int:
    parser = argparse.ArgumentParser(description="Export a PowerPoint deck and rebuild its flashcard metadata.")
    parser.add_argument("deck_folder", type=Path, help="Folder such as decks/fish_morphology")
    parser.add_argument("--name", help="Display name. Defaults to the PowerPoint filename.")
    parser.add_argument("--skip-export", action="store_true", help="Only rebuild JSON from existing slide PNG files.")
    args = parser.parse_args()

    project_root = Path(__file__).resolve().parent.parent
    deck_folder = args.deck_folder.resolve()
    powerpoint = find_powerpoint(deck_folder)
    titles = read_titles(powerpoint)

    if not args.skip_export:
        export_slides(powerpoint, deck_folder, Path(__file__).resolve().parent)

    images = [deck_folder / f"slide-{number}.png" for number in range(1, len(titles) + 1)]
    missing = [image.name for image in images if not image.exists()]
    if missing:
        raise FileNotFoundError("Missing exported slide images: " + ", ".join(missing))

    deck_id = slugify(deck_folder.name)
    display_name = args.name or powerpoint.stem.replace("_", " ").strip().title()
    cards = [
        {"id": f"slide-{number}", "image": f"slide-{number}.png", "answer": title}
        for number, title in enumerate(titles, start=1)
    ]
    deck_data = {"id": deck_id, "name": display_name, "cards": cards}
    (deck_folder / "deck.json").write_text(json.dumps(deck_data, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    rebuild_index(project_root)
    print(f"Updated {display_name}: {len(cards)} cards.")
    return 0


if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except Exception as error:
        print(f"ERROR: {error}", file=sys.stderr)
        raise SystemExit(1)
